from __future__ import annotations

import io
import json
import math
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt

from .models import Blueprint


def export_json_bytes(blueprint: Blueprint) -> bytes:
    return json.dumps(blueprint.model_dump(), indent=2).encode("utf-8")


def _dxf_line(x1: float, y1: float, x2: float, y2: float, layer: str = "0") -> list[str]:
    return [
        "0",
        "LINE",
        "8",
        layer,
        "10",
        f"{x1:.4f}",
        "20",
        f"{y1:.4f}",
        "30",
        "0.0",
        "11",
        f"{x2:.4f}",
        "21",
        f"{y2:.4f}",
        "31",
        "0.0",
    ]


def _dxf_circle(cx: float, cy: float, r: float, layer: str = "0") -> list[str]:
    return [
        "0",
        "CIRCLE",
        "8",
        layer,
        "10",
        f"{cx:.4f}",
        "20",
        f"{cy:.4f}",
        "30",
        "0.0",
        "40",
        f"{r:.4f}",
    ]


def _dxf_text(x: float, y: float, text: str, height: float = 3.2, layer: str = "ANNOT") -> list[str]:
    return [
        "0",
        "TEXT",
        "8",
        layer,
        "10",
        f"{x:.4f}",
        "20",
        f"{y:.4f}",
        "30",
        "0.0",
        "40",
        f"{height:.4f}",
        "1",
        text,
    ]


def export_dxf_bytes(blueprint: Blueprint) -> bytes:
    d = blueprint.dimensions

    body_h = d.body_height_mm
    overall_h = d.overall_height_mm

    r_bottom = d.body_bottom_diameter_mm * 0.5
    r_max = d.body_max_diameter_mm * 0.5
    r_neck = d.neck_diameter_mm * 0.5
    r_head = d.head_top_diameter_mm * 0.5

    right_profile = [
        (r_bottom, 0.0),
        (r_max, body_h * 0.30),
        (r_max * 0.98, body_h * 0.68),
        (r_neck, body_h),
        (r_neck * 1.18, body_h + (overall_h - body_h) * 0.45),
        (r_head, overall_h),
    ]

    left_profile = [(-x, y) for x, y in right_profile]

    top_center_x = 280.0
    top_center_y = overall_h * 0.60

    entities: list[str] = []

    for i in range(len(right_profile) - 1):
        x1, y1 = right_profile[i]
        x2, y2 = right_profile[i + 1]
        entities.extend(_dxf_line(x1, y1, x2, y2, "SIDE"))

    for i in range(len(left_profile) - 1):
        x1, y1 = left_profile[i]
        x2, y2 = left_profile[i + 1]
        entities.extend(_dxf_line(x1, y1, x2, y2, "SIDE"))

    entities.extend(_dxf_line(0, -8, 0, overall_h + 12, "CENTER"))
    entities.extend(_dxf_line(-r_bottom, 0, r_bottom, 0, "SIDE"))
    entities.extend(_dxf_line(-r_head, overall_h, r_head, overall_h, "SIDE"))

    entities.extend(_dxf_circle(top_center_x, top_center_y, r_head, "TOP"))
    entities.extend(_dxf_circle(top_center_x, top_center_y, r_neck, "TOP"))
    entities.extend(_dxf_circle(top_center_x, top_center_y, d.insert_outer_diameter_mm * 0.5, "TOP"))
    entities.extend(_dxf_circle(top_center_x, top_center_y, d.insert_inner_diameter_mm * 0.5, "TOP"))

    entities.extend(_dxf_text(-95, overall_h + 14, f"Overall Height: {overall_h:.1f} mm"))
    entities.extend(_dxf_text(-95, overall_h + 9, f"Body Max Dia: {d.body_max_diameter_mm:.1f} mm"))
    entities.extend(_dxf_text(-95, overall_h + 4, f"Head Top Dia: {d.head_top_diameter_mm:.1f} mm"))
    entities.extend(_dxf_text(top_center_x - 50, top_center_y - r_head - 8, "Top View"))

    sections = [
        "0",
        "SECTION",
        "2",
        "HEADER",
        "0",
        "ENDSEC",
        "0",
        "SECTION",
        "2",
        "TABLES",
        "0",
        "ENDSEC",
        "0",
        "SECTION",
        "2",
        "ENTITIES",
        *entities,
        "0",
        "ENDSEC",
        "0",
        "EOF",
    ]

    return ("\n".join(sections) + "\n").encode("ascii", errors="ignore")


@dataclass
class MeshBuilder:
    vertices: list[tuple[float, float, float]] = field(default_factory=list)
    faces: list[tuple[int, int, int]] = field(default_factory=list)

    def add_vertex(self, x: float, y: float, z: float) -> int:
        self.vertices.append((x, y, z))
        return len(self.vertices)

    def add_face(self, a: int, b: int, c: int) -> None:
        self.faces.append((a, b, c))

    def add_lathe(
        self,
        profile: list[tuple[float, float]],
        segments: int = 56,
        close_bottom: bool = False,
        close_top: bool = False,
    ) -> None:
        if len(profile) < 2:
            return

        rings: list[list[int]] = []
        for r, y in profile:
            ring: list[int] = []
            for i in range(segments):
                theta = (2.0 * math.pi * i) / segments
                x = r * math.cos(theta)
                z = r * math.sin(theta)
                ring.append(self.add_vertex(x, y, z))
            rings.append(ring)

        for j in range(len(rings) - 1):
            cur = rings[j]
            nxt = rings[j + 1]
            for i in range(segments):
                i2 = (i + 1) % segments
                a = cur[i]
                b = nxt[i]
                c = nxt[i2]
                d = cur[i2]
                self.add_face(a, b, c)
                self.add_face(a, c, d)

        if close_bottom:
            center = self.add_vertex(0.0, profile[0][1], 0.0)
            bottom_ring = rings[0]
            for i in range(segments):
                i2 = (i + 1) % segments
                self.add_face(center, bottom_ring[i2], bottom_ring[i])

        if close_top:
            center = self.add_vertex(0.0, profile[-1][1], 0.0)
            top_ring = rings[-1]
            for i in range(segments):
                i2 = (i + 1) % segments
                self.add_face(center, top_ring[i], top_ring[i2])

    def add_cylinder(
        self,
        radius: float,
        height: float,
        y0: float,
        segments: int = 40,
        close_bottom: bool = False,
        close_top: bool = False,
    ) -> None:
        self.add_lathe(
            [(radius, y0), (radius, y0 + height)],
            segments=segments,
            close_bottom=close_bottom,
            close_top=close_top,
        )

    def add_tube_path(
        self,
        points: list[tuple[float, float, float]],
        radius: float,
        radial_segments: int = 14,
    ) -> None:
        if len(points) < 2:
            return

        rings: list[list[int]] = []
        for idx, point in enumerate(points):
            x, y, z = point
            if idx == 0:
                tx = points[idx + 1][0] - x
                ty = points[idx + 1][1] - y
                tz = points[idx + 1][2] - z
            elif idx == len(points) - 1:
                tx = x - points[idx - 1][0]
                ty = y - points[idx - 1][1]
                tz = z - points[idx - 1][2]
            else:
                tx = points[idx + 1][0] - points[idx - 1][0]
                ty = points[idx + 1][1] - points[idx - 1][1]
                tz = points[idx + 1][2] - points[idx - 1][2]

            t_len = math.sqrt(tx * tx + ty * ty + tz * tz) or 1.0
            tx, ty, tz = tx / t_len, ty / t_len, tz / t_len

            # Build an orthonormal frame around the tangent.
            ux, uy, uz = (0.0, 0.0, 1.0)
            if abs(tz) > 0.92:
                ux, uy, uz = (0.0, 1.0, 0.0)

            nx = ty * uz - tz * uy
            ny = tz * ux - tx * uz
            nz = tx * uy - ty * ux
            n_len = math.sqrt(nx * nx + ny * ny + nz * nz) or 1.0
            nx, ny, nz = nx / n_len, ny / n_len, nz / n_len

            bx = ty * nz - tz * ny
            by = tz * nx - tx * nz
            bz = tx * ny - ty * nx

            ring: list[int] = []
            for i in range(radial_segments):
                theta = (2.0 * math.pi * i) / radial_segments
                rx = radius * (math.cos(theta) * nx + math.sin(theta) * bx)
                ry = radius * (math.cos(theta) * ny + math.sin(theta) * by)
                rz = radius * (math.cos(theta) * nz + math.sin(theta) * bz)
                ring.append(self.add_vertex(x + rx, y + ry, z + rz))
            rings.append(ring)

        for j in range(len(rings) - 1):
            cur = rings[j]
            nxt = rings[j + 1]
            for i in range(radial_segments):
                i2 = (i + 1) % radial_segments
                a = cur[i]
                b = nxt[i]
                c = nxt[i2]
                d = cur[i2]
                self.add_face(a, b, c)
                self.add_face(a, c, d)


def export_obj_bytes(blueprint: Blueprint) -> bytes:
    d = blueprint.dimensions
    mesh = MeshBuilder()

    body_h = d.body_height_mm
    head_start = body_h - d.head_neck_overlap_mm
    overall_h = d.overall_height_mm

    r_bottom = d.body_bottom_diameter_mm * 0.5
    r_max = d.body_max_diameter_mm * 0.5
    r_neck = d.neck_diameter_mm * 0.5
    r_head = d.head_top_diameter_mm * 0.5

    body_profile = [
        (r_bottom, 0.0),
        (r_max, body_h * 0.30),
        (r_max * 0.98, body_h * 0.68),
        (r_neck, body_h),
    ]

    head_profile = [
        (r_neck, head_start),
        (r_neck * 1.18, head_start + (overall_h - head_start) * 0.45),
        (r_head, overall_h),
    ]

    mesh.add_lathe(body_profile, segments=64, close_bottom=True, close_top=False)
    mesh.add_lathe(head_profile, segments=64, close_bottom=False, close_top=False)

    # Insert/filter collar
    insert_h = d.insert_height_mm
    insert_outer = d.insert_outer_diameter_mm * 0.5
    insert_inner = d.insert_inner_diameter_mm * 0.5
    insert_y0 = overall_h - insert_h

    mesh.add_cylinder(insert_outer, insert_h, insert_y0, segments=44, close_bottom=False, close_top=False)
    mesh.add_cylinder(insert_inner, insert_h, insert_y0, segments=44, close_bottom=False, close_top=False)

    # Simple annular top face for insert
    ring_profile = [
        (insert_inner, overall_h),
        (insert_outer, overall_h),
    ]
    mesh.add_lathe(ring_profile, segments=44, close_bottom=False, close_top=False)

    # Handle as tube path
    anchor_x = r_head + d.handle_offset_mm
    handle_points = [
        (anchor_x, overall_h * 0.84, 0.0),
        (anchor_x + d.handle_length_mm * 0.45, overall_h * 0.72, 0.0),
        (anchor_x + d.handle_length_mm * 0.38, overall_h * 0.46, 0.0),
        (anchor_x + d.handle_length_mm * 0.12, overall_h * 0.34 - d.handle_drop_mm * 0.08, 0.0),
    ]
    mesh.add_tube_path(handle_points, radius=max(d.handle_thickness_mm * 0.5, 1.5), radial_segments=14)

    lines = ["# Curved-head teapot OBJ export", "o teapot"]
    for x, y, z in mesh.vertices:
        lines.append(f"v {x:.6f} {y:.6f} {z:.6f}")
    for a, b, c in mesh.faces:
        lines.append(f"f {a} {b} {c}")

    return ("\n".join(lines) + "\n").encode("ascii", errors="ignore")


def export_pptx_bytes(blueprint: Blueprint) -> bytes:
    d = blueprint.dimensions

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9.0), Inches(0.7))
    title_tf = title.text_frame
    title_tf.text = blueprint.title
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True

    summary = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(6.3), Inches(2.2))
    summary_tf = summary.text_frame
    summary_tf.word_wrap = True
    summary_tf.text = (
        f"Target: {d.cups_target:.1f} cups ({d.capacity_target_ml:.0f} ml)\n"
        f"Estimated Capacity: {d.estimated_capacity_ml:.0f} ml\n"
        f"Overall Height: {d.overall_height_mm:.1f} mm ({d.overall_height_mm / 25.4:.2f} in)\n"
        f"Body Max Dia: {d.body_max_diameter_mm:.1f} mm ({d.body_max_diameter_mm / 25.4:.2f} in)\n"
        f"Head Top Dia: {d.head_top_diameter_mm:.1f} mm ({d.head_top_diameter_mm / 25.4:.2f} in)\n"
        f"Material Baseline: Stainless Steel 304 + nylon handle + silicone gasket"
    )

    notes_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.9), Inches(9.2), Inches(2.6))
    notes_tf = notes_box.text_frame
    notes_tf.word_wrap = True
    notes_tf.text = "Design Notes"
    notes_tf.paragraphs[0].font.bold = True
    for note in blueprint.analysis_notes[:6]:
        p = notes_tf.add_paragraph()
        p.text = f"- {note}"
        p.level = 0

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    t2 = slide2.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9.0), Inches(0.6))
    t2.text_frame.text = "Material Suggestions"
    t2.text_frame.paragraphs[0].font.size = Pt(24)
    t2.text_frame.paragraphs[0].font.bold = True

    rows = max(len(blueprint.materials) + 1, 2)
    cols = 5
    table_shape = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(9.2), Inches(5.8))
    table = table_shape.table
    headers = ["Part", "Selected Material", "Recommended", "Confidence", "Notes"]
    for col, text in enumerate(headers):
        table.cell(0, col).text = text

    for idx, mat in enumerate(blueprint.materials, start=1):
        table.cell(idx, 0).text = mat.part_name
        table.cell(idx, 1).text = mat.selected or mat.recommended
        table.cell(idx, 2).text = mat.recommended
        table.cell(idx, 3).text = f"{mat.confidence * 100:.0f}%"
        table.cell(idx, 4).text = mat.notes

    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    t3 = slide3.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9.0), Inches(0.6))
    t3.text_frame.text = "Manufacturing Bill of Materials"
    t3.text_frame.paragraphs[0].font.size = Pt(24)
    t3.text_frame.paragraphs[0].font.bold = True

    bom_rows = max(len(blueprint.bom) + 1, 2)
    bom_cols = 6
    bom_shape = slide3.shapes.add_table(bom_rows, bom_cols, Inches(0.4), Inches(1.0), Inches(9.4), Inches(5.9))
    bom_table = bom_shape.table
    bom_headers = ["Part", "Material", "Process", "Thk (mm)", "Qty", "Mass (g)"]
    for col, text in enumerate(bom_headers):
        bom_table.cell(0, col).text = text

    for idx, item in enumerate(blueprint.bom, start=1):
        bom_table.cell(idx, 0).text = item.part_name
        bom_table.cell(idx, 1).text = item.material
        bom_table.cell(idx, 2).text = item.process
        bom_table.cell(idx, 3).text = f"{item.thickness_mm:.2f}"
        bom_table.cell(idx, 4).text = str(item.quantity)
        bom_table.cell(idx, 5).text = f"{item.mass_estimate_g:.1f}"

    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()
